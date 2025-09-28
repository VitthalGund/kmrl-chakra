import fitz  # PyMuPDF
import cv2
import numpy as np
from PIL import Image
import io
import pypandoc
import tempfile
import os
import shutil
import pytesseract
from dotenv import load_dotenv
import google.generativeai as genai
import docx
from pptx import Presentation

from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.responses import JSONResponse

# --- Initial Setup ---
load_dotenv()
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

app = FastAPI(
    title="Universal Document Processor API",
    description="A comprehensive API to convert, enhance, separate, OCR, and summarize any document."
)


# --- Core Processing Functions ---

def apply_enhancement_pipeline(image: Image.Image) -> Image.Image:
    """Applies a robust pipeline to enhance image quality for OCR."""
    open_cv_image = cv2.cvtColor(np.array(image.convert("RGB")), cv2.COLOR_RGB2BGR)
    gray = cv2.cvtColor(open_cv_image, cv2.COLOR_BGR2GRAY)
    denoised = cv2.fastNlMeansDenoising(gray, None, h=10, templateWindowSize=7, searchWindowSize=21)
    sharpen_kernel = np.array([[-1, -1, -1], [-1, 9, -1], [-1, -1, -1]])
    sharpened = cv2.filter2D(denoised, -1, sharpen_kernel)
    binary = cv2.adaptiveThreshold(
        sharpened, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
        cv2.THRESH_BINARY, 11, 2
    )
    return Image.fromarray(binary)

def extract_text_from_image(enhanced_image: Image.Image) -> str:
    """Performs OCR on an enhanced PIL Image."""
    try:
        config = r'--oem 3 --psm 6'
        return pytesseract.image_to_string(enhanced_image, config=config, lang='eng').strip()
    except Exception as e:
        return f"[ERROR] OCR failed: {e}"

def generate_gemini_summary(enhanced_image: Image.Image, ocr_text: str) -> str:
    """Generates a summary using Gemini, combining visual and text analysis."""
    if not os.getenv("GOOGLE_API_KEY"):
        return "[INFO] Google API key not configured. Skipping summary."
        
    model = genai.GenerativeModel('gemini-1.5-flash')
    try:
        prompt = f"""
        Analyze the provided image and its OCR text to generate a concise summary.
        Focus on the document's purpose, key data points, and any important entities.

        OCR Text:
        ---
        {ocr_text}
        ---
        """
        response = model.generate_content([prompt, enhanced_image])
        return response.text.strip() if response.text else "[ERROR] Gemini returned an empty response."
    except Exception as e:
        return f"[ERROR] Gemini API call failed: {e}"

def analyze_image_data(image_bytes: bytes, page_num: int, img_index: int, temp_dir: str):
    """Full pipeline for a single image: enhance, OCR, summarize."""
    pil_image = Image.open(io.BytesIO(image_bytes))
    enhanced_image = apply_enhancement_pipeline(pil_image)
    ocr_text = extract_text_from_image(enhanced_image)
    summary = generate_gemini_summary(enhanced_image, ocr_text)
    
    img_filename = f"enhanced_page_{page_num + 1}_img_{img_index + 1}.png"
    enhanced_image.save(os.path.join(temp_dir, img_filename))
    
    return {
        "source_page": page_num + 1,
        "filename": img_filename,
        "ocr_text": ocr_text,
        "summary": summary
    }, enhanced_image

def create_output_pdfs(text_content: str, enhanced_images: list, temp_dir: str):
    """Creates two separate PDFs: one for text and one for images."""
    # Create Text-Only PDF using Pandoc for robust, multi-page conversion
    text_pdf_path = os.path.join(temp_dir, "text_only_output.pdf")
    try:
        pypandoc.convert_text(text_content, 'pdf', format='markdown',
                              outputfile=text_pdf_path, extra_args=['--pdf-engine=pdflatex'])
    except Exception:
        # Fallback to a simpler text insertion if LaTeX isn't installed
        text_doc = fitz.open()
        page = text_doc.new_page()
        page.insert_textbox(page.rect, text_content, fontsize=11, fontname="helv")
        text_doc.save(text_pdf_path)
        text_doc.close()

    # Create Image-Only PDF
    image_pdf_path = os.path.join(temp_dir, "images_only_output.pdf")
    if enhanced_images:
        img_docs = [Image.open(os.path.join(temp_dir, res['filename'])) for res in enhanced_images]
        img_docs[0].save(
            image_pdf_path, "PDF", resolution=100.0,
            save_all=True, append_images=img_docs[1:]
        )
    else:
        fitz.open().save(image_pdf_path) # Create an empty PDF if no images were found

    return text_pdf_path, image_pdf_path

# --- Document Specific Handlers ---

def process_pdf_file(file_bytes: bytes, temp_dir: str):
    pdf_doc = fitz.open(stream=file_bytes, filetype="pdf")
    all_text = ""
    image_analysis_results = []

    for page_num, page in enumerate(pdf_doc):
        all_text += page.get_text("text") + "\n\n"
        for img_index, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base_image = pdf_doc.extract_image(xref)
            image_bytes = base_image["image"]
            analysis, _ = analyze_image_data(image_bytes, page_num, img_index, temp_dir)
            image_analysis_results.append(analysis)

    return all_text, image_analysis_results

def process_docx_file(file_bytes: bytes, temp_dir: str):
    doc = docx.Document(io.BytesIO(file_bytes))
    all_text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
    image_analysis_results = []
    
    for i, rel in enumerate(doc.part.rels.values()):
        if "image" in rel.target_ref:
            image_bytes = rel.target_part.blob
            analysis, _ = analyze_image_data(image_bytes, 0, i, temp_dir)
            image_analysis_results.append(analysis)
            
    return all_text, image_analysis_results

def process_pptx_file(file_bytes: bytes, temp_dir: str):
    prs = Presentation(io.BytesIO(file_bytes))
    all_text = ""
    image_analysis_results = []
    img_index = 0

    for slide_num, slide in enumerate(prs.slides):
        all_text += f"\n--- Slide {slide_num + 1} ---\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                all_text += shape.text + "\n"
            if hasattr(shape, "image"):
                image_bytes = shape.image.blob
                analysis, _ = analyze_image_data(image_bytes, slide_num, img_index, temp_dir)
                image_analysis_results.append(analysis)
                img_index += 1
                
    return all_text, image_analysis_results

# --- Main API Endpoint ---

@app.post("/process-document/")
async def process_document(file: UploadFile = File(...)):
    temp_dir = "temp_processing"
    if os.path.exists(temp_dir):
        shutil.rmtree(temp_dir)
    os.makedirs(temp_dir)

    try:
        file_bytes = await file.read()
        original_filename = file.filename
        file_extension = os.path.splitext(original_filename)[1].lower()

        all_text = ""
        image_analysis_results = []

        if file_extension == ".pdf":
            all_text, image_analysis_results = process_pdf_file(file_bytes, temp_dir)
        elif file_extension == ".docx":
            all_text, image_analysis_results = process_docx_file(file_bytes, temp_dir)
        elif file_extension == ".pptx":
            all_text, image_analysis_results = process_pptx_file(file_bytes, temp_dir)
        elif file.content_type and file.content_type.startswith("image/"):
            analysis, _ = analyze_image_data(file_bytes, 0, 0, temp_dir)
            image_analysis_results.append(analysis)
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported file type: {original_filename}. Please upload a PDF, DOCX, PPTX, or image file.")

        text_pdf_path, image_pdf_path = create_output_pdfs(all_text, image_analysis_results, temp_dir)
        
        return JSONResponse({
            "status": "success",
            "original_filename": original_filename,
            "images_found": len(image_analysis_results),
            "image_analysis": image_analysis_results,
            "text_only_pdf_path": text_pdf_path,
            "images_only_pdf_path": image_pdf_path,
        })

    except Exception as e:
        if isinstance(e, HTTPException):
            raise e
        # Provide a more detailed error log for debugging
        print(f"An unexpected error occurred: {str(e)}")
        raise HTTPException(status_code=500, detail=f"An unexpected error occurred processing {original_filename}. Details: {str(e)}")

